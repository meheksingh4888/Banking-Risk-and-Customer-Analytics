"""
Create a PowerPoint presentation for Banking Risk Customer Analytics project
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    """Create Banking Risk Analytics presentation"""
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define color scheme
    BLUE = RGBColor(0, 102, 204)
    DARK_GRAY = RGBColor(51, 51, 51)
    LIGHT_GRAY = RGBColor(242, 242, 242)
    
    def add_title_slide(title, subtitle):
        """Add a title slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Add background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BLUE
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_p = title_frame.paragraphs[0]
        title_p.text = title
        title_p.font.size = Pt(54)
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(255, 255, 255)
        title_p.alignment = PP_ALIGN.CENTER
        
        # Add subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.word_wrap = True
        subtitle_p = subtitle_frame.paragraphs[0]
        subtitle_p.text = subtitle
        subtitle_p.font.size = Pt(24)
        subtitle_p.font.color.rgb = RGBColor(255, 255, 255)
        subtitle_p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def add_content_slide(title, content_list):
        """Add a content slide with bullet points"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Add white background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)
        
        # Add title bar
        title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
        title_shape.fill.solid()
        title_shape.fill.fore_color.rgb = BLUE
        title_shape.line.color.rgb = BLUE
        
        # Add title text
        title_frame = title_shape.text_frame
        title_p = title_frame.paragraphs[0]
        title_p.text = title
        title_p.font.size = Pt(40)
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(255, 255, 255)
        title_p.space_before = Pt(10)
        title_p.space_after = Pt(10)
        
        # Add content
        left = Inches(0.75)
        top = Inches(1.5)
        width = Inches(8.5)
        height = Inches(5.5)
        
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        
        for i, item in enumerate(content_list):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = item
            p.level = 0
            p.font.size = Pt(18)
            p.font.color.rgb = DARK_GRAY
            p.space_after = Pt(10)
        
        return slide
    
    # Slide 1: Title Slide
    add_title_slide(
        "Banking Risk Customer Analytics",
        "End-to-End Data Analytics Project\nPrepared by Mehek Singh"
    )
    
    # Slide 2: Project Overview
    add_content_slide(
        "Project Overview",
        [
            "✓ Complete end-to-end data analytics solution",
            "✓ Simulates real-world banking operations",
            "✓ Demonstrates database design, ETL, and analytics",
            "✓ Interactive Power BI dashboards for insights",
            "✓ Professional documentation and automation"
        ]
    )
    
    # Slide 3: Business Problem
    add_content_slide(
        "Business Challenges",
        [
            "• Identifying high-risk customers",
            "• Monitoring loan defaults",
            "• Tracking branch performance",
            "• Detecting fraudulent transactions",
            "• Understanding customer behavior",
            "• Need for centralized reporting system"
        ]
    )
    
    # Slide 4: Technology Stack
    add_content_slide(
        "Technology Stack",
        [
            "• Database: PostgreSQL",
            "• Data Processing: Python, Pandas, SQLAlchemy",
            "• Analytics: SQL (Basic to Advanced)",
            "• Visualization: Power BI",
            "• Version Control: Git & GitHub",
            "• Automation: Batch scripts & Python ETL"
        ]
    )
    
    # Slide 5: Database Design
    add_content_slide(
        "Database Architecture",
        [
            "• 8 Main Tables: Customers, Accounts, Transactions, Loans, Branches, Merchants, Loan Types, Account Types",
            "• Normalized schema with Primary & Foreign Keys",
            "• Referential Integrity enforced",
            "• Performance optimized with Indexes",
            "• Views and Materialized Views for reporting"
        ]
    )
    
    # Slide 6: ETL Process
    add_content_slide(
        "ETL Pipeline",
        [
            "1. Extract: Raw CSV datasets from multiple sources",
            "2. Clean: Data validation & standardization",
            "3. Transform: Column mapping & type conversion",
            "4. Load: PostgreSQL database ingestion",
            "5. Verify: Record counts & data quality checks",
            "6. Report: Automated validation summaries"
        ]
    )
    
    # Slide 7: SQL Analytics
    add_content_slide(
        "SQL Analysis Coverage",
        [
            "• Basic SQL Queries",
            "• Intermediate & Advanced SQL",
            "• Window Functions & CTEs",
            "• Stored Procedures & Functions",
            "• Triggers & Transactions",
            "• Query Optimization & Indexing"
        ]
    )
    
    # Slide 8: Python Analytics
    add_content_slide(
        "Python Analysis Modules",
        [
            "✓ Customer Analysis - Behavior & demographics",
            "✓ Account Analysis - Balance trends",
            "✓ Transaction Analysis - Volume & patterns",
            "✓ Loan Analysis - Portfolio distribution",
            "✓ Branch Analysis - Performance metrics",
            "✓ Risk Analysis - High-risk identification"
        ]
    )
    
    # Slide 9: Power BI Dashboards
    add_content_slide(
        "Interactive Dashboards",
        [
            "• Executive Dashboard - KPIs & trends",
            "• Customer Dashboard - Segmentation & behavior",
            "• Transaction Dashboard - Channels & volumes",
            "• Loan Dashboard - Performance & defaults",
            "• Branch Dashboard - Rankings & metrics",
            "• Risk Dashboard - Fraud & anomalies"
        ]
    )
    
    # Slide 10: Key Insights
    add_content_slide(
        "Business Insights",
        [
            "• Customer distribution across states/cities",
            "• Account balance trends over time",
            "• Monthly transaction volume patterns",
            "• Loan portfolio distribution",
            "• Top performing branches and merchants",
            "• High-risk customer identification"
        ]
    )
    
    # Slide 11: Challenges & Solutions
    add_content_slide(
        "Challenges Overcome",
        [
            "• Schema design complexity → Normalized design",
            "• Data validation issues → Automated validation",
            "• Query optimization → Indexing & CTEs",
            "• ETL debugging → Comprehensive logging",
            "• Relationship modeling → Power BI best practices"
        ]
    )
    
    # Slide 12: Future Roadmap
    add_content_slide(
        "Future Enhancements",
        [
            "• Real-time data streaming",
            "• ML-based fraud detection",
            "• Credit risk prediction models",
            "• Customer churn prediction",
            "• Cloud deployment (AWS/Azure)",
            "• REST API integration"
        ]
    )
    
    # Slide 13: Conclusion
    add_content_slide(
        "Conclusion",
        [
            "✓ Complete end-to-end analytics workflow",
            "✓ Real-world enterprise practices",
            "✓ Scalable & maintainable architecture",
            "✓ Strong portfolio demonstration",
            "✓ Ready for Data Analyst, BI Developer, & Data Engineer roles"
        ]
    )
    
    # Slide 14: Thank You
    add_title_slide(
        "Thank You",
        "Questions & Discussion\nMehek Singh | Banking Risk Analytics"
    )
    
    # Save presentation
    output_path = r"c:\Users\sakas\SQL PROJECTS\📁 Banking-Risk-Customer-Analytics\Presentation\Banking_Risk_Analytics_Generated.pptx"
    prs.save(output_path)
    print(f"✓ Presentation created: {output_path}")
    print(f"✓ Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    create_presentation()
