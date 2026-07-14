from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor

# -------------------------------------------------------
# Create Presentation
# -------------------------------------------------------
prs = Presentation()

# Widescreen Layout
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# Theme Colors
BLUE = RGBColor(0, 70, 140)
LIGHT_BLUE = RGBColor(220, 235, 250)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)
GRAY = RGBColor(90, 90, 90)


# -------------------------------------------------------
# Helper Functions
# -------------------------------------------------------

def add_title(slide, title_text):
    """Add a professional slide title."""
    title_box = slide.shapes.add_textbox(
        Inches(0.4),
        Inches(0.2),
        Inches(12),
        Inches(0.5)
    )

    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = BLUE


def add_footer(slide):
    """Footer and slide number."""
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(7.1),
        Inches(13.33),
        Inches(0.05)
    )

    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.color.rgb = BLUE

    footer = slide.shapes.add_textbox(
        Inches(0.3),
        Inches(7.15),
        Inches(4),
        Inches(0.3)
    )

    p = footer.text_frame.paragraphs[0]
    p.text = "Prepared By: Mehek Singh"
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY

    number = slide.shapes.add_textbox(
        Inches(12.5),
        Inches(7.1),
        Inches(0.5),
        Inches(0.3)
    )

    p = number.text_frame.paragraphs[0]
    p.text = str(len(prs.slides))
    p.font.size = Pt(10)
    p.alignment = PP_ALIGN.RIGHT


def add_bullets(slide, items):
    """Add bullet list."""
    textbox = slide.shapes.add_textbox(
        Inches(0.8),
        Inches(1.0),
        Inches(11.5),
        Inches(5.5)
    )

    tf = textbox.text_frame

    for i, item in enumerate(items):

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = item
        p.level = 0
        p.font.size = Pt(22)
        p.font.color.rgb = BLACK
        p.space_after = Pt(10)


# -------------------------------------------------------
# Slide 1 - Title
# -------------------------------------------------------

slide = prs.slides.add_slide(prs.slide_layouts[6])

title = slide.shapes.add_textbox(
    Inches(0.8),
    Inches(1),
    Inches(11),
    Inches(1)
)

p = title.text_frame.paragraphs[0]
p.text = "Banking Risk Customer Analytics"
p.font.size = Pt(34)
p.font.bold = True
p.font.color.rgb = BLUE

p = title.text_frame.add_paragraph()
p.text = "End-to-End Data Analytics Project"
p.font.size = Pt(24)

body = slide.shapes.add_textbox(
    Inches(1),
    Inches(3),
    Inches(8),
    Inches(3)
)

tf = body.text_frame

p = tf.paragraphs[0]
p.text = "Prepared By"
p.font.size = Pt(22)
p.font.bold = True

p = tf.add_paragraph()
p.text = "Mehek Singh"
p.font.size = Pt(24)

p = tf.add_paragraph()
p.text = ""
p.font.size = Pt(8)

p = tf.add_paragraph()
p.text = "Technology"
p.font.bold = True
p.font.size = Pt(22)

technologies = [
    "PostgreSQL",
    "SQL",
    "Python",
    "Power BI",
    "Pandas"
]

for tech in technologies:
    p = tf.add_paragraph()
    p.text = f"• {tech}"
    p.font.size = Pt(20)

add_footer(slide)


# -------------------------------------------------------
# Slide 2 - Problem Statement
# -------------------------------------------------------

slide = prs.slides.add_slide(prs.slide_layouts[6])

add_title(slide, "Problem Statement")

items = [
    "Problem",
    "Banks generate huge volumes of customer, account, loan and transaction data.",
    "",
    "Without analytics it is difficult to:",
    "• Detect fraud",
    "• Monitor customer risk",
    "• Track loan defaults",
    "• Analyze branch performance",
    "• Understand customer behavior",
    "",
    "Goal",
    "Build an end-to-end analytics solution to help management make data-driven decisions."
]

add_bullets(slide, items)

add_footer(slide)


# -------------------------------------------------------
# Slide 3 - Business Requirements
# -------------------------------------------------------

slide = prs.slides.add_slide(prs.slide_layouts[6])

add_title(slide, "Business Requirements")

items = [
    "Objectives",
    "",
    "✓ Customer Analytics",
    "✓ Account Analysis",
    "✓ Transaction Monitoring",
    "✓ Loan Portfolio Analysis",
    "✓ Fraud Detection",
    "✓ Branch Performance",
    "✓ Executive Dashboard"
]

add_bullets(slide, items)

add_footer(slide)


# -------------------------------------------------------
# Slide 4 - Project Architecture
# -------------------------------------------------------

slide = prs.slides.add_slide(prs.slide_layouts[6])

add_title(slide, "Project Architecture")

architecture = [
    "CSV Files",
    "↓",
    "Python ETL",
    "↓",
    "PostgreSQL Database",
    "↓",
    "SQL Analysis",
    "↓",
    "Python EDA",
    "↓",
    "Power BI Dashboard",
    "↓",
    "Business Insights"
]

box = slide.shapes.add_textbox(
    Inches(4),
    Inches(1),
    Inches(5),
    Inches(5)
)

tf = box.text_frame

for i, step in enumerate(architecture):

    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()

    p.text = step
    p.font.size = Pt(24)
    p.alignment = PP_ALIGN.CENTER

add_footer(slide)

# -------------------------------------------------------
# Slide 5 - Database Design
# -------------------------------------------------------

slide = prs.slides.add_slide(prs.slide_layouts[6])

add_title(slide, "Database Design")

items = [
    "Database Tables",
    "",
    "• Customers",
    "• Accounts",
    "• Transactions",
    "• Loans",
    "• Branches",
    "• Merchants",
    "• Fraud Alerts",
    "• Account Types",
    "• Loan Types",
    "",
    "Database Concepts",
    "",
    "• Primary Keys",
    "• Foreign Keys",
    "• Relationships",
    "• Normalization",
    "• Referential Integrity"
]

add_bullets(slide, items)

add_footer(slide)


# -------------------------------------------------------
# Slide 6 - ER Diagram
# -------------------------------------------------------

slide = prs.slides.add_slide(prs.slide_layouts[6])

add_title(slide, "ER Diagram")

box = slide.shapes.add_textbox(
    Inches(1),
    Inches(1.5),
    Inches(11),
    Inches(4)
)

tf = box.text_frame

p = tf.paragraphs[0]
p.text = "Entity Relationship Diagram"
p.font.size = Pt(28)
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

p = tf.add_paragraph()
p.text = ""
p.font.size = Pt(12)

p = tf.add_paragraph()
p.text = "Customers"
p.font.size = Pt(22)
p.alignment = PP_ALIGN.CENTER

p = tf.add_paragraph()
p.text = "│"
p.font.size = Pt(22)
p.alignment = PP_ALIGN.CENTER

p = tf.add_paragraph()
p.text = "Accounts"
p.font.size = Pt(22)
p.alignment = PP_ALIGN.CENTER

p = tf.add_paragraph()
p.text = "│"
p.font.size = Pt(22)
p.alignment = PP_ALIGN.CENTER

p = tf.add_paragraph()
p.text = "Transactions"
p.font.size = Pt(22)
p.alignment = PP_ALIGN.CENTER

p = tf.add_paragraph()
p.text = "│"
p.font.size = Pt(22)
p.alignment = PP_ALIGN.CENTER

p = tf.add_paragraph()
p.text = "Loans"
p.font.size = Pt(22)
p.alignment = PP_ALIGN.CENTER

add_footer(slide)


# -------------------------------------------------------
# Slide 7 - SQL Work
# -------------------------------------------------------

slide = prs.slides.add_slide(prs.slide_layouts[6])

add_title(slide, "SQL Work")

items = [
    "SQL Concepts Used",
    "",
    "• Joins",
    "• Common Table Expressions (CTEs)",
    "• Window Functions",
    "• Views",
    "• Stored Procedures",
    "• Functions",
    "• Indexing",
    "• Query Optimization",
    "• Aggregations",
    "• GROUP BY",
    "• HAVING"
]

add_bullets(slide, items)

query_box = slide.shapes.add_textbox(
    Inches(7),
    Inches(1.2),
    Inches(5.5),
    Inches(4)
)

tf = query_box.text_frame

p = tf.paragraphs[0]
p.text = "Example SQL Query"
p.font.bold = True
p.font.size = Pt(18)

queries = [
    "SELECT customer_id,",
    "COUNT(*) AS loan_count,",
    "SUM(loan_amount)",
    "FROM loans",
    "GROUP BY customer_id;",
]

for q in queries:
    para = tf.add_paragraph()
    para.text = q
    para.font.name = "Courier New"
    para.font.size = Pt(16)

add_footer(slide)


# -------------------------------------------------------
# Slide 8 - Python ETL
# -------------------------------------------------------

slide = prs.slides.add_slide(prs.slide_layouts[6])

add_title(slide, "Python ETL")

items = [
    "ETL Process",
    "",
    "• Reading CSV Files",
    "• Data Cleaning",
    "• Missing Value Handling",
    "• Duplicate Removal",
    "• Data Validation",
    "• Data Type Conversion",
    "• Loading Data into PostgreSQL",
    "• Error Handling",
    "• Logging",
    "",
    "Python Libraries Used",
    "",
    "• pandas",
    "• psycopg2",
    "• numpy"
]

add_bullets(slide, items)

add_footer(slide)

# -------------------------------------------------------
# Slide 9 - Python Analysis
# -------------------------------------------------------

slide = prs.slides.add_slide(prs.slide_layouts[6])

add_title(slide, "Python Analysis")

items = [
    "Python Analysis Modules",
    "",
    "• customer_analysis.py",
    "• account_analysis.py",
    "• transaction_analysis.py",
    "• loan_analysis.py",
    "• branch_analysis.py",
    "• merchant_analysis.py",
    "• fraud_analysis.py",
    "• executive_summary.py",
    "",
    "Analysis Performed",
    "",
    "• Customer segmentation",
    "• Account performance",
    "• Transaction trends",
    "• Loan analysis",
    "• Branch-wise KPIs",
    "• Fraud detection"
]

add_bullets(slide, items)

add_footer(slide)


# -------------------------------------------------------
# Slide 10 - Power BI Dashboard
# -------------------------------------------------------

slide = prs.slides.add_slide(prs.slide_layouts[6])

add_title(slide, "Power BI Dashboard")

items = [
    "Executive Dashboard Features",
    "",
    "• KPI Cards",
    "• Interactive Charts",
    "• Filters & Slicers",
    "• Customer Analysis",
    "• Branch Performance",
    "• Transaction Monitoring",
    "• Loan Portfolio",
    "• Risk Analysis",
    "",
    "Business users can interact with the dashboard",
    "to explore insights and monitor banking performance."
]

add_bullets(slide, items)

add_footer(slide)


# -------------------------------------------------------
# Slide 11 - Business Insights
# -------------------------------------------------------

slide = prs.slides.add_slide(prs.slide_layouts[6])

add_title(slide, "Business Insights")

items = [
    "Project Summary",
    "",
    "• Total Customers : 20,000",
    "• Total Accounts : 30,000",
    "• Total Transactions : 500,000",
    "• Loan Portfolio : ₹25.27 Billion",
    "• Failed Transactions : 9,863",
    "• International Transactions : 10,040",
    "• Loan Defaults : 788",
    "",
    "These insights help management make",
    "data-driven business decisions."
]

add_bullets(slide, items)

add_footer(slide)


# -------------------------------------------------------
# Slide 12 - Risk Analysis
# -------------------------------------------------------

slide = prs.slides.add_slide(prs.slide_layouts[6])

add_title(slide, "Risk Analysis")

items = [
    "Major Risks Identified",
    "",
    "• Dormant Accounts",
    "• Failed Transactions",
    "• High Loan Exposure",
    "• International Transactions",
    "• Low Credit Score Customers",
    "",
    "Risk Monitoring Benefits",
    "",
    "• Early fraud detection",
    "• Better customer monitoring",
    "• Improved loan recovery",
    "• Reduced financial risk",
    "• Better compliance"
]

add_bullets(slide, items)

add_footer(slide)

# -------------------------------------------------------
# Slide 13 - Recommendations
# -------------------------------------------------------

slide = prs.slides.add_slide(prs.slide_layouts[6])

add_title(slide, "Recommendations")

items = [
    "Recommended Actions",
    "",
    "• Monitor high-risk customers regularly.",
    "• Reduce failed transactions through system optimization.",
    "• Improve branch performance using KPI monitoring.",
    "• Promote digital banking adoption.",
    "• Enhance loan approval strategy with risk scoring.",
    "• Strengthen fraud monitoring and alert systems.",
    "",
    "Expected Benefits",
    "",
    "• Better customer satisfaction",
    "• Reduced operational risk",
    "• Improved profitability",
    "• Faster decision-making"
]

add_bullets(slide, items)

add_footer(slide)


# -------------------------------------------------------
# Slide 14 - Future Scope
# -------------------------------------------------------

slide = prs.slides.add_slide(prs.slide_layouts[6])

add_title(slide, "Future Scope")

items = [
    "Future Enhancements",
    "",
    "• Real-time Dashboard",
    "• Machine Learning Risk Prediction",
    "• AI-based Fraud Detection",
    "• Customer Churn Prediction",
    "• Cloud Deployment",
    "• Power BI Service Integration",
    "• Automated Email Alerts",
    "• Predictive Analytics",
    "• Mobile Dashboard Access"
]

add_bullets(slide, items)

add_footer(slide)


# -------------------------------------------------------
# Slide 15 - Thank You
# -------------------------------------------------------

slide = prs.slides.add_slide(prs.slide_layouts[6])

title_box = slide.shapes.add_textbox(
    Inches(2),
    Inches(1.3),
    Inches(9),
    Inches(1)
)

tf = title_box.text_frame

p = tf.paragraphs[0]
p.text = "THANK YOU"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = BLUE
p.alignment = PP_ALIGN.CENTER

p = tf.add_paragraph()
p.text = ""
p.font.size = Pt(10)

p = tf.add_paragraph()
p.text = "Questions?"
p.font.size = Pt(26)
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

p = tf.add_paragraph()
p.text = ""
p.font.size = Pt(10)

p = tf.add_paragraph()
p.text = "Mehek Singh"
p.font.size = Pt(24)
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

p = tf.add_paragraph()
p.text = "Data Analyst | SQL | Python | Power BI"
p.font.size = Pt(20)
p.alignment = PP_ALIGN.CENTER

add_footer(slide)


# -------------------------------------------------------
# Save Presentation
# -------------------------------------------------------

output_file = "Banking_Risk_Presentation.pptx"

prs.save(output_file)

print("=" * 50)
print("Presentation created successfully!")
print(f"File Name : {output_file}")
print(f"Total Slides : {len(prs.slides)}")
print("=" * 50)